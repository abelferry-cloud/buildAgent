"""PPT generation skill using python-pptx."""

import os
from typing import Optional

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    Presentation = None

from agent.tools.base import Tool, ToolResult


class PptGenerateTool(Tool):
    """Generate a PowerPoint presentation."""

    name = "ppt_generate"
    description = "Generate a PowerPoint presentation with slides."

    def __init__(self, output_dir: str = "output"):
        super().__init__()
        self.output_dir = os.path.join(os.getcwd(), output_dir)
        os.makedirs(self.output_dir, exist_ok=True)

    def execute(
        self,
        title: str,
        slides: list[dict],
        output_filename: str = "presentation.pptx",
    ) -> ToolResult:
        """Generate a PPTX file."""
        if Presentation is None:
            return ToolResult(
                output="",
                error="python-pptx not installed. Run: pip install python-pptx"
            )

        try:
            prs = Presentation()

            # Title slide
            if slides:
                first = slides[0]
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                slide.shapes.title.text = first.get("title", title)
                if "content" in first:
                    slide.placeholders[1].text = first["content"]

            # Content slides
            for slide_data in slides[1:]:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = slide_data.get("title", "")
                if "content" in slide_data:
                    slide.placeholders[1].text = slide_data["content"]

            output_path = os.path.join(self.output_dir, output_filename)
            prs.save(output_path)
            return ToolResult(output=f"PPT 已生成: {output_path}")
        except Exception as e:
            return ToolResult(output="", error=str(e))
